from pptx import Presentation

def process(file_path: str) -> str:
    """
    Processes a presentation file (.pptx).
    """
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error reading PPTX {file_path}: {e}")
        return ""
